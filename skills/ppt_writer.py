from typing import Dict, Any, List
from core.prompt_engine import PromptEngine
from core.base_skill import BaseSkill
from pptx import Presentation
from pptx.util import Inches
import os
import json
from config.logger import logger

class PPTWriter(BaseSkill):
    skill_name = "ppt_writer"
    
    def __init__(self):
        super().__init__()
        self.prompt_engine = PromptEngine()
        
    async def validate_input(self, input_data: Dict[str, Any]) -> bool:
        if "content" not in input_data or not input_data["content"]:
            logger.error("내용이 없습니다")
            return False
        return True
        
    async def _process_internal(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        내용을 PPT로 변환
        
        Args:
            input_data: 변환할 내용이 포함된 입력 데이터
            
        Returns:
            생성된 PPT 파일 경로와 메타데이터
        """
        content = input_data["content"]
        
        # 내용 구조화
        prompt = f"""
        다음 내용을 PPT 슬라이드 형식으로 구조화해주세요.
        각 슬라이드는 다음 JSON 형식으로 반환해주세요:
        
        {{
            "slides": [
                {{
                    "title": "슬라이드 제목",
                    "bullets": ["항목 1", "항목 2", ...],
                    "notes": "발표자 노트"
                }}
            ]
        }}
        
        내용:
        {content}
        """
        
        structured_content = await self.prompt_engine.run_prompt(
            prompt=prompt,
            temperature=0.7
        )
        
        # JSON 파싱
        slides_data = json.loads(structured_content)
        
        # PPT 생성
        prs = Presentation()
        for slide in slides_data["slides"]:
            # 슬라이드 추가
            layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
            current_slide = prs.slides.add_slide(layout)
            
            # 제목 설정
            if current_slide.shapes.title:
                current_slide.shapes.title.text = slide["title"]
            
            # 내용 추가
            if current_slide.shapes.placeholders[1]:
                content_shape = current_slide.shapes.placeholders[1]
                tf = content_shape.text_frame
                
                for bullet in slide["bullets"]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            
            # 노트 추가 (있는 경우)
            if "notes" in slide and slide["notes"]:
                notes_slide = current_slide.notes_slide
                notes_slide.notes_text_frame.text = slide["notes"]
        
        # 파일 저장
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        file_path = os.path.join(output_dir, f"presentation_{hash(content)}.pptx")
        prs.save(file_path)
        
        return {
            "file_path": file_path,
            "slide_count": len(slides_data["slides"]),
            "metadata": {
                "titles": [slide["title"] for slide in slides_data["slides"]],
                "total_bullets": sum(len(slide["bullets"]) for slide in slides_data["slides"])
            }
        }
